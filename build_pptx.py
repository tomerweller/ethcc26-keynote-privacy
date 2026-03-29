#!/usr/bin/env python3
"""Build the EthCC Privacy slide deck as a native PPTX with animations."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import textwrap

# --- Brand Colors ---
BLACK = RGBColor(0x0F, 0x0F, 0x0F)
OFFWHITE = RGBColor(0xF6, 0xF7, 0xF8)
WARMGRAY = RGBColor(0xD6, 0xD2, 0xC4)
YELLOW = RGBColor(0xFD, 0xDA, 0x24)
LAVENDER = RGBColor(0xB7, 0xAC, 0xE8)

# --- Slide dimensions (16:9) ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# --- Layout constants ---
MARGIN_LEFT = Inches(1.2)
MARGIN_TOP_TITLE = Inches(2.0)
CONTENT_WIDTH = Inches(10.5)
TITLE_HEIGHT = Inches(1.5)
BULLET_START_Y = Inches(3.8)
BULLET_SPACING = Inches(0.65)


def set_slide_bg(slide, color=BLACK):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_appear_animation(slide, shape, delay_idx=0):
    """Add 'Appear on click' animation to a shape."""
    # Build the animation XML
    timing = slide.element.find(qn('p:transition'))

    # Get or create timing node
    sld = slide.element
    timing_node = sld.find(qn('p:timing'))
    if timing_node is None:
        timing_node = deepcopy(sld.makeelement(qn('p:timing'), {}))
        sld.append(timing_node)

    tn_lst = timing_node.find(qn('p:tnLst'))
    if tn_lst is None:
        tn_lst = timing_node.makeelement(qn('p:tnLst'), {})
        timing_node.append(tn_lst)

    par_seq = tn_lst.find(qn('p:par'))
    if par_seq is None:
        # Create the root parallel container
        par_seq = tn_lst.makeelement(qn('p:par'), {})
        tn_lst.append(par_seq)

        c_tn_root = par_seq.makeelement(qn('p:cTn'), {
            'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'
        })
        par_seq.append(c_tn_root)

        child_tn_lst = c_tn_root.makeelement(qn('p:childTnLst'), {})
        c_tn_root.append(child_tn_lst)

        seq = child_tn_lst.makeelement(qn('p:seq'), {
            'concurrent': '1', 'nextAc': 'seek'
        })
        child_tn_lst.append(seq)

        seq_ctn = seq.makeelement(qn('p:cTn'), {
            'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'
        })
        seq.append(seq_ctn)

        seq_child_lst = seq_ctn.makeelement(qn('p:childTnLst'), {})
        seq_ctn.append(seq_child_lst)

        # prev/next conditions for seq
        prev_cond = seq.makeelement(qn('p:prevCondLst'), {})
        seq.append(prev_cond)
        cond_prev = prev_cond.makeelement(qn('p:cond'), {'evt': 'onPrev', 'delay': '0'})
        prev_cond.append(cond_prev)
        tgt_el_prev = cond_prev.makeelement(qn('p:tgtEl'), {})
        cond_prev.append(tgt_el_prev)
        sld_tgt_prev = tgt_el_prev.makeelement(qn('p:sldTgt'), {})
        tgt_el_prev.append(sld_tgt_prev)

        next_cond = seq.makeelement(qn('p:nextCondLst'), {})
        seq.append(next_cond)
        cond_next = next_cond.makeelement(qn('p:cond'), {'evt': 'onNext', 'delay': '0'})
        next_cond.append(cond_next)
        tgt_el_next = cond_next.makeelement(qn('p:tgtEl'), {})
        cond_next.append(tgt_el_next)
        sld_tgt_next = tgt_el_next.makeelement(qn('p:sldTgt'), {})
        tgt_el_next.append(sld_tgt_next)
    else:
        c_tn_root = par_seq.find(qn('p:cTn'))
        child_tn_lst = c_tn_root.find(qn('p:childTnLst'))
        seq = child_tn_lst.find(qn('p:seq'))
        seq_ctn = seq.find(qn('p:cTn'))
        seq_child_lst = seq_ctn.find(qn('p:childTnLst'))

    # Get next available id
    existing_ids = [int(el.get('id', '0')) for el in sld.iter(qn('p:cTn'))]
    next_id = max(existing_ids) + 1 if existing_ids else 3

    # Each click-triggered animation is a par inside the seq childTnLst
    click_par = seq_child_lst.makeelement(qn('p:par'), {})
    seq_child_lst.append(click_par)

    click_ctn = click_par.makeelement(qn('p:cTn'), {
        'id': str(next_id), 'fill': 'hold'
    })
    click_par.append(click_ctn)

    # Start condition: on click
    stCondLst = click_ctn.makeelement(qn('p:stCondLst'), {})
    click_ctn.append(stCondLst)
    cond = stCondLst.makeelement(qn('p:cond'), {'delay': '0'})
    stCondLst.append(cond)

    click_child_lst = click_ctn.makeelement(qn('p:childTnLst'), {})
    click_ctn.append(click_child_lst)

    inner_par = click_child_lst.makeelement(qn('p:par'), {})
    click_child_lst.append(inner_par)

    inner_ctn = inner_par.makeelement(qn('p:cTn'), {
        'id': str(next_id + 1), 'presetID': '1', 'presetClass': 'entr',
        'presetSubtype': '0', 'fill': 'hold', 'nodeType': 'clickEffect'
    })
    inner_par.append(inner_ctn)

    inner_stcond = inner_ctn.makeelement(qn('p:stCondLst'), {})
    inner_ctn.append(inner_stcond)
    inner_cond = inner_stcond.makeelement(qn('p:cond'), {'delay': '0'})
    inner_stcond.append(inner_cond)

    inner_child = inner_ctn.makeelement(qn('p:childTnLst'), {})
    inner_ctn.append(inner_child)

    # The set element that makes the shape visible
    set_el = inner_child.makeelement(qn('p:set'), {})
    inner_child.append(set_el)

    set_ctn = set_el.makeelement(qn('p:cBhvr'), {})
    set_el.append(set_ctn)

    set_ctn_inner = set_ctn.makeelement(qn('p:cTn'), {
        'id': str(next_id + 2), 'dur': '1', 'fill': 'hold'
    })
    set_ctn.append(set_ctn_inner)

    set_stcond = set_ctn_inner.makeelement(qn('p:stCondLst'), {})
    set_ctn_inner.append(set_stcond)
    set_cond = set_stcond.makeelement(qn('p:cond'), {'delay': '0'})
    set_stcond.append(set_cond)

    tgt_el = set_ctn.makeelement(qn('p:tgtEl'), {})
    set_ctn.append(tgt_el)

    sp_tgt = tgt_el.makeelement(qn('p:spTgt'), {'spid': str(shape.shape_id)})
    tgt_el.append(sp_tgt)

    to_el = set_el.makeelement(qn('p:to'), {})
    set_el.append(to_el)
    str_val = to_el.makeelement(qn('p:strVal'), {'val': 'visible'})
    to_el.append(str_val)


def make_shape_invisible(shape):
    """Set initial visibility to hidden (for animation)."""
    sp = shape.element
    # Add or modify cNvPr to not be visible initially -- we do this via
    # setting the shape to hidden in the spPr
    sp_pr = sp.find(qn('p:spPr'))
    if sp_pr is None:
        sp_pr = sp.find(qn('p:grpSpPr'))


def add_title_slide(prs, title, subtitle=None, notes=""):
    """Add a title slide with large centered title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Title
    top = Inches(2.2) if subtitle else Inches(2.8)
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.name = "Lora"
    p.font.color.rgb = OFFWHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        txBox2 = slide.shapes.add_textbox(MARGIN_LEFT, top + Inches(2.2), CONTENT_WIDTH, Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.name = "Inter"
        p2.font.color.rgb = WARMGRAY
        p2.alignment = PP_ALIGN.LEFT

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_title_only_slide(prs, title, notes=""):
    """Add a slide with just a centered title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(2.8), CONTENT_WIDTH, Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.name = "Lora"
    p.font.color.rgb = OFFWHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_bullets_slide(prs, title, bullets, notes="", animate=True):
    """Add a slide with title and bullet points that appear on click.

    bullets: list of tuples (label, description) or just strings.
    If tuple, label is yellow bold, description is warm gray.
    If string, entire line is yellow bold.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Inches(1.5), CONTENT_WIDTH, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.name = "Lora"
    p.font.color.rgb = OFFWHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Each bullet as its own textbox (for individual animation)
    bullet_shapes = []
    for i, bullet in enumerate(bullets):
        top = Inches(3.4) + Inches(0.7) * i
        bx = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, Inches(0.6))
        btf = bx.text_frame
        btf.word_wrap = True
        btf.auto_size = None
        bp = btf.paragraphs[0]

        if isinstance(bullet, tuple):
            label, desc = bullet
            run1 = bp.add_run()
            run1.text = label
            run1.font.size = Pt(22)
            run1.font.name = "Inter"
            run1.font.color.rgb = YELLOW
            run1.font.bold = True

            run2 = bp.add_run()
            run2.text = f" -- {desc}"
            run2.font.size = Pt(22)
            run2.font.name = "Inter"
            run2.font.color.rgb = WARMGRAY
            run2.font.bold = False
        else:
            run1 = bp.add_run()
            run1.text = bullet
            run1.font.size = Pt(22)
            run1.font.name = "Inter"
            run1.font.color.rgb = YELLOW
            run1.font.bold = True

        bp.alignment = PP_ALIGN.LEFT
        bullet_shapes.append(bx)

    # Add appear-on-click animations
    if animate:
        for i, shape in enumerate(bullet_shapes):
            add_appear_animation(slide, shape, i)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


# ============================================================
# BUILD THE DECK
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# --- Slide 1: Title ---
add_title_slide(prs,
    "Privacy That Doesn't\nBreak the Chain",
    subtitle="Tomer Weller · Stellar.org",
    notes="Hi, I'm Tomer Weller, CPO at the Stellar Development Foundation. I've been in crypto for almost 9 years and for the last couple of years, privacy has been my main focus. This talk is about what we've learned in discussions with institutions about privacy."
)

# --- Slide 2: Why Privacy? ---
add_title_only_slide(prs, "Why Privacy?",
    notes="""If you're in this room you're probably already sold on privacy. But just to level set:

Blockchain is the only financial system where every transaction is visible to everyone. Your paycheck is private. Your bank balance is private. But the moment any of that moves on-chain, it's public.

For individuals, this means anyone with a block explorer can reconstruct your financial life. For institutions, it's worse -- volumes, relationships, trading strategies, payroll. This is competitive intelligence, exposed to every competitor, in real time.

When we talk to financial institutions, they make it clear: not having a way to manage visibility is a non-starter.

Transition: Easy -- we've been working on privacy for 10 years, no? Not exactly."""
)

# --- Slide 3: The Privacy Mistake ---
add_title_only_slide(prs, "The Privacy Mistake",
    notes="""Yeah, we've been working on privacy for a while but with a strong focus on cypherpunk values -- complete financial freedom.

And that's awesome. But for the next wave of adoption we need institutions -- we need them to feel comfortable bringing their assets on chain, making payroll, accepting merchant payments, engaging with blockchain in a meaningful way.

For these institutions privacy means something different -- it means managing visibility in a way that allows them to stay competitive but also in compliance.

Transition: The mistake is that we've built absolute privacy. But what institutions are telling us is that privacy is a spectrum."""
)

# --- Slide 4: Privacy is a Spectrum ---
add_title_only_slide(prs, "Privacy is a Spectrum",
    notes="""Privacy is not binary. It's a spectrum with an enormous number of permutations:

Jurisdictions vary. A regulated asset issued in the US adheres to a completely different regulatory regime than a similar one issued in Germany.

Use cases vary. Payroll needs different privacy than trading, which needs different privacy than remittances.

Technology varies. ZK proofs, FHE, MPC, TEEs -- each has different tradeoffs in performance, trust assumptions, and what they can actually hide.

The real problem is that we've been treating privacy as a single feature when it's actually a multi-dimensional design space.

Transition: This is what I call "Privacy Dead Ends"."""
)

# --- Slide 5: Privacy Dead Ends ---
add_bullets_slide(prs, "Privacy Dead Ends", [
    ("All-or-Nothing", "full transparency or full privacy, nothing in between"),
    ("Permissioned Chains", "private by default, but you lose everything blockchain offers"),
],
    notes="""When we started looking at privacy for Stellar, we realized that most privacy solutions end up in one of two failure modes, or "dead ends".

All-or-Nothing: Privacy solutions that have two modes: full transparency or full privacy, nothing in between. This can take the shape of a smart contract pool, an L2 or an appchain - but the premise is the same. Unfortunately, this doesn't leave room for regulated financial institutions to engage with.

Permissioned Chains: I've been in crypto almost 9 years and it's very depressing - as soon as you think that the concept of permissioned chains is dead, they come back to life. Recently they've come back to life on the pretense of being private. Of course they're private - they're a fucking database on your computer, bro. Just to be clear: institutions are looking for open-participation networks - they see the value of interoperability, composability and verifiability. Permissioned chains are a dead end.

Transition: So how do you avoid these privacy dead-ends?"""
)

# --- Slide 6: Transparency First ---
add_bullets_slide(prs, "Transparency First", [
    ("Open, transparent base layer", "maximally auditable source of truth"),
    ("Privacy protocols built on top", "configurable, opt-in, different tech for different needs"),
],
    notes="""The answer is simple: a clear architectural separation.

An open, transparent base layer -- a maximally auditable source of truth. This is the default - this is the primary issuance platform. Anyone can verify - interoperability is trivial.

Privacy protocols built on top. Different privacy guarantees, different administrative controls, different underlying tech.

Why transparency first? This might come as a surprise given that this is a talk about privacy but transparency is a core value proposition of blockchains. The fact that an executive at an asset manager company can go on a block explorer and see the total supply of an asset and its distribution - that is a feature - not a bug - they love that shit.

With that said, our base layers should provide a rich set of building blocks for privacy. Stellar's X-Ray upgrade is an example of this approach -- adding ZK primitives (BN254, Poseidon) at the protocol level so that application developers can build configurable privacy without compromising the transparency of the base chain.

Transition: We keep talking about different privacy protocols for different use cases. Let's demonstrate with some concrete examples."""
)

# --- Slide 7: Private Payments ---
add_bullets_slide(prs, "Private Payments", [
    ("Confidential Tokens", "hide the what (amounts, balances)"),
    ("Privacy Pools", "hide the who (sender, receiver)"),
],
    notes="""We're going to focus on Private Payments. There's interesting stuff happening in private defi but these are still a ways out.

The two main families of private payment protocols that we see are confidential tokens and privacy pools.

Confidential Tokens: Hide the what -- payment amounts and balances are concealed. But the sender and receiver are still visible on-chain. You know who's transacting, you just can't see how much. This is great for when counterparty relationships are known - for example with payroll, people know I'm employed by SDF but my salary is private. There are a bunch of different implementations of this. Stellar is part of the confidential token association and we're working with OpenZeppelin on an implementation. Narrower privacy guarantees but the tech is fairly scalable and it's a bit easier for compliance.

Privacy Pools: Payment protocols that also hide the who -- funds are mixed so that sender and receivers identities are hidden. This provides anonymity, not just confidentiality. Much stronger privacy guarantees but they inherently mix funds, which means legitimate funds can be commingled with illicit funds - which there are ways to tackle.

Transition: This demonstrates different privacy guarantees, and the question is how do we build administrative controls to enable compliance with these solutions and others."""
)

# --- Slide 8: Privacy Compliance* Menu ---
add_bullets_slide(prs, "Privacy Compliance* Menu", [
    "Selective Disclosure (View Keys)",
    "Non-Selective Disclosure (Auditor Keys)",
    "Association Sets",
    "Forced Transparent Withdrawals",
    "Clawback",
],
    notes="""At Stellar, we've been working with various builders in the space to define a compliance menu of opt-in administrative controls. The asterisk is deliberate. First of all I'm not a lawyer, and also regulators haven't defined what compliance looks like for on-chain privacy yet. These are the administrative controls we're starting to see as requirements, this is not a comprehensive list and the idea is to have these as configurable opt-in.

Selective Disclosure (View Keys): The user chooses to reveal specific transaction details to specific parties. This allows the user to show a source of funds on demand.

Non-Selective Disclosure (Auditor Keys): A third party authority holds a key that can view all transactions within a scope. This means that if a law enforcement agency has a subpoena they have an actual door to knock on.

Association Sets: An allow list controlled by a pool operator - it ensures that all funds mixed come from approved addresses and reduces the risk of commingling with illicit funds.

Forced Transparent Withdrawals: What happens if an account gets revoked? They're forced to withdraw transparently -- no hiding behind the pool on the way out. 0xBow calls this rage quit.

Clawback: If clawbacks are triggering for you then you are in the wrong room and you are not ready for what's coming next. Tokenization is at full speed and regulated assets often require clawback capabilities.

Transition: These are dials, not switches. Different combinations serve different jurisdictions, use cases, and risk profiles."""
)

# --- Slide 9: Closing ---
add_title_slide(prs,
    "Privacy That Doesn't\nBreak the Chain",
    notes="""So to wrap up:

We've been building privacy with ideology first and that's great -- but it's not enough. Institutions need privacy too, and for them it's more nuanced.

Privacy is a spectrum. Different jurisdictions, different use cases, different tech. There's no single answer and that's okay -- that's actually the point.

The right approach is transparency first -- keep the base layer open, build configurable privacy on top. Confidential tokens, privacy pools, a compliance menu of administrative controls. Dials, not switches.

All the building blocks are in place - we have a non-hostile administration in the US, we have institutions that are at the table and see the value in blockchain, and the tech is ready.

There are no excuses - we're building privacy that doesn't break the chain. If that's interesting to you -- come build with us on Stellar.

Thanks."""
)

# Save
prs.save('slides.pptx')
print(f"Created slides.pptx: {len(prs.slides)} slides with animations")
